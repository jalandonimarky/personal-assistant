"""Document builders. Invoked by the MCP server with a JSON spec on stdin.

Kept in Python because openpyxl and python-pptx already exist on the machine and
are the correct tools; reimplementing OOXML by hand would be worse in every way.
The server passes structured data, never a command string.
"""
import json, subprocess, sys, os, html


def spreadsheet(spec, out):
    from openpyxl import Workbook
    from openpyxl.styles import Font
    wb = Workbook()
    wb.remove(wb.active)
    for sheet in spec["sheets"]:
        ws = wb.create_sheet(title=str(sheet.get("name", "Sheet"))[:31])
        rows = sheet.get("rows", [])
        for r in rows:
            ws.append(["" if c is None else c for c in r])
        if rows and sheet.get("header", True):
            for cell in ws[1]:
                cell.font = Font(bold=True)
            ws.freeze_panes = "A2"
        # Width from content, capped so one long cell cannot blow out the sheet.
        for col in ws.columns:
            width = max((len(str(c.value)) for c in col if c.value is not None), default=8)
            ws.column_dimensions[col[0].column_letter].width = min(max(width + 2, 9), 60)
    wb.save(out)
    return {"sheets": len(spec["sheets"])}


def presentation(spec, out):
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    for i, slide in enumerate(spec["slides"]):
        bullets = slide.get("bullets") or []
        # Title-only layout when there is nothing to bullet.
        layout = prs.slide_layouts[0 if (i == 0 and not bullets) else 1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = str(slide.get("title", ""))
        if bullets and len(s.placeholders) > 1:
            tf = s.placeholders[1].text_frame
            tf.text = str(bullets[0])
            for b in bullets[1:]:
                para = tf.add_paragraph()
                para.text = str(b)
                para.level = 0
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(18)
        elif not bullets and len(s.placeholders) > 1 and slide.get("subtitle"):
            s.placeholders[1].text = str(slide["subtitle"])
        if slide.get("notes"):
            s.notes_slide.notes_text_frame.text = str(slide["notes"])
    prs.save(out)
    return {"slides": len(spec["slides"])}


def document(spec, out):
    """HTML → .docx via LibreOffice. python-docx is not installed, and soffice
    is already here and handles styling more predictably than hand-built XML."""
    parts = [
        "<html><head><meta charset='utf-8'><style>",
        "body{font-family:Helvetica,Arial,sans-serif;font-size:11pt;line-height:1.5}",
        "h1{font-size:20pt}h2{font-size:14pt}table{border-collapse:collapse;width:100%}",
        "td,th{border:1px solid #999;padding:5px;text-align:left}",
        "</style></head><body>",
    ]
    if spec.get("title"):
        parts.append(f"<h1>{html.escape(str(spec['title']))}</h1>")
    for b in spec.get("blocks", []):
        kind, text = b.get("type", "p"), html.escape(str(b.get("text", "")))
        if kind == "heading":
            parts.append(f"<h2>{text}</h2>")
        elif kind == "bullets":
            items = "".join(f"<li>{html.escape(str(x))}</li>" for x in b.get("items", []))
            parts.append(f"<ul>{items}</ul>")
        elif kind == "table":
            rows = b.get("rows", [])
            head = "".join(f"<th>{html.escape(str(c))}</th>" for c in rows[0]) if rows else ""
            body = "".join(
                "<tr>" + "".join(f"<td>{html.escape(str(c))}</td>" for c in r) + "</tr>"
                for r in rows[1:]
            )
            parts.append(f"<table><tr>{head}</tr>{body}</table>")
        else:
            parts.append(f"<p>{text}</p>")
    parts.append("</body></html>")

    # A single-dot basename: soffice derives the output name by swapping the
    # extension, and "name.docx.src.html" makes it produce "name.docx.src.docx"
    # which it then refuses to export.
    outdir = os.path.dirname(out) or "."
    stem = "_docbuild_" + str(os.getpid())
    tmp = os.path.join(outdir, stem + ".html")
    with open(tmp, "w") as f:
        f.write("".join(parts))
    r = subprocess.run(
        # HTML opens as a Writer/Web document, which has no default docx export
        # filter — the filter has to be named explicitly or soffice aborts.
        ["soffice", "--headless", "--convert-to", "docx:MS Word 2007 XML",
         "--outdir", outdir, tmp],
        capture_output=True, timeout=180,
    )
    produced = os.path.join(outdir, stem + ".docx")
    if not os.path.exists(produced):
        raise RuntimeError(f"soffice produced nothing: {r.stderr.decode()[:200]}")
    os.replace(produced, out)
    os.remove(tmp)
    return {"blocks": len(spec.get("blocks", []))}


def read_any(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    out.append(" | ".join("" if c is None else str(c) for c in row))
        return "\n".join(out)
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for i, s in enumerate(prs.slides, 1):
            out.append(f"# Slide {i}")
            for shape in s.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    if ext in (".docx", ".odt", ".rtf", ".pdf"):
        # soffice converts to text without needing a reader library per format.
        outdir = os.path.dirname(path) or "."
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", outdir, path],
            capture_output=True, timeout=180,
        )
        txt = os.path.join(outdir, os.path.splitext(os.path.basename(path))[0] + ".txt")
        if os.path.exists(txt):
            with open(txt, errors="replace") as f:
                body = f.read()
            os.remove(txt)
            return body
        raise RuntimeError("could not extract text")
    with open(path, errors="replace") as f:
        return f.read()


if __name__ == "__main__":
    req = json.load(sys.stdin)
    op, out = req["op"], req.get("out")
    try:
        if op == "spreadsheet":
            print(json.dumps({"ok": True, **spreadsheet(req["spec"], out)}))
        elif op == "presentation":
            print(json.dumps({"ok": True, **presentation(req["spec"], out)}))
        elif op == "document":
            print(json.dumps({"ok": True, **document(req["spec"], out)}))
        elif op == "read":
            print(json.dumps({"ok": True, "text": read_any(req["path"])[:60000]}))
        else:
            print(json.dumps({"ok": False, "error": f"unknown op {op}"}))
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"{type(e).__name__}: {e}"}))
